import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pydantic import BaseModel, Field
from typing import List, Optional

# 定義 Pydantic 模型，用於 Gemini 的結構化輸出 (Structured Outputs)
class SlideContent(BaseModel):
    subtitle: Optional[str] = Field(None, description="投影片副標題 (主要用於 title_slide 首頁)")
    points: Optional[List[str]] = Field(None, description="條列重點 (標準條列式 bullets 版面使用)")
    column1_title: Optional[str] = Field(None, description="左側欄位標題 (雙欄 two_columns 版面使用)")
    column1_points: Optional[List[str]] = Field(None, description="左側欄位重點 (雙欄 two_columns 版面使用)")
    column2_title: Optional[str] = Field(None, description="右側欄位標題 (雙欄 two_columns 版面使用)")
    column2_points: Optional[List[str]] = Field(None, description="右側欄位重點 (雙欄 two_columns 版面使用)")

class SlideOutline(BaseModel):
    title: str = Field(..., description="投影片單頁標題")
    layout: str = Field(..., description="版面類型：'title_slide' (首頁), 'bullets' (標準條列), 'two_columns' (左右雙欄)")
    content: SlideContent = Field(..., description="投影片單頁內容")

class PPTOutline(BaseModel):
    title: str = Field(..., description="整份簡報的總標題")
    slides: List[SlideOutline] = Field(..., description="投影片頁面清單")

# 簡報色彩主題設定
THEMES = {
    "light": {
        "bg": RGBColor(248, 249, 250),          # 淺灰白
        "primary": RGBColor(26, 26, 26),        # 近乎黑
        "secondary": RGBColor(74, 74, 74),      # 深灰
        "accent": RGBColor(0, 86, 179)          # 科技藍
    },
    "dark": {
        "bg": RGBColor(17, 24, 39),             # 藍黑色 (slate-900)
        "primary": RGBColor(249, 250, 251),     # 亮灰白
        "secondary": RGBColor(209, 213, 219),   # 淺灰
        "accent": RGBColor(245, 158, 11)        # 琥珀橘
    }
}

# 字型名稱 (支援 Windows 常用微軟正黑體，macOS 會自動映射或匹配)
FONT_NAME = "Microsoft JhengHei"

def create_presentation(outline: PPTOutline, theme_name: str = "light") -> Presentation:
    """根據大綱資料結構與配色主題建立簡報並回傳 Presentation 物件。"""
    prs = Presentation()
    
    # 設定 16:9 寬螢幕尺寸 (13.333 x 7.5 英吋)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    theme = THEMES.get(theme_name.lower(), THEMES["light"])
    blank_layout = prs.slide_layouts[6]  # 取得空白投影片版型
    
    for i, slide_data in enumerate(outline.slides):
        slide = prs.slides.add_slide(blank_layout)
        
        # 套用背景顏色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = theme["bg"]
        
        layout_type = slide_data.layout
        if layout_type == "title_slide":
            add_title_slide(slide, slide_data.title, slide_data.content.subtitle, theme)
        elif layout_type == "two_columns":
            add_two_columns_slide(slide, slide_data.title, slide_data.content, theme)
        else:  # 預設為 bullets 條列式
            add_bullets_slide(slide, slide_data.title, slide_data.content.points, theme)
            
    return prs

def add_title_slide(slide, title: str, subtitle: Optional[str], theme: dict):
    """加入精緻的首頁投影片。"""
    # 左側修飾的彩色粗線條
    shape = slide.shapes.add_shape(
        1,  # msoShapeRectangle
        Inches(1.0), Inches(2.2), Inches(0.15), Inches(3.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["accent"]
    shape.line.fill.background()  # 無框線
    
    # 建立主標與副標文字框
    tx_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.5), Inches(3.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # 主標題
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(46)
    p_title.font.bold = True
    p_title.font.color.rgb = theme["primary"]
    p_title.space_after = Pt(20)
    
    # 副標題
    if subtitle:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = FONT_NAME
        p_sub.font.size = Pt(24)
        p_sub.font.color.rgb = theme["secondary"]

def add_bullets_slide(slide, title: str, points: Optional[List[str]], theme: dict):
    """加入標準條列式投影片。"""
    add_header(slide, title, theme)
    
    # 建立內文文字框
    tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    if points:
        for idx, pt in enumerate(points):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = f"•  {pt}"
            p.font.name = FONT_NAME
            p.font.size = Pt(20)
            p.font.color.rgb = theme["secondary"]
            p.space_after = Pt(14)

def add_two_columns_slide(slide, title: str, content: SlideContent, theme: dict):
    """加入左右雙欄對比投影片。"""
    add_header(slide, title, theme)
    
    # 左側欄位
    tx_box_l = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.5))
    tf_l = tx_box_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
    
    if content.column1_title:
        p_title_l = tf_l.paragraphs[0]
        p_title_l.text = content.column1_title
        p_title_l.font.name = FONT_NAME
        p_title_l.font.size = Pt(22)
        p_title_l.font.bold = True
        p_title_l.font.color.rgb = theme["accent"]
        p_title_l.space_after = Pt(12)
        
    if content.column1_points:
        for idx, pt in enumerate(content.column1_points):
            p = tf_l.add_paragraph() if (idx > 0 or content.column1_title) else tf_l.paragraphs[0]
            p.text = f"•  {pt}"
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = theme["secondary"]
            p.space_after = Pt(10)

    # 右側欄位
    tx_box_r = slide.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.3), Inches(4.5))
    tf_r = tx_box_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
    
    if content.column2_title:
        p_title_r = tf_r.paragraphs[0]
        p_title_r.text = content.column2_title
        p_title_r.font.name = FONT_NAME
        p_title_r.font.size = Pt(22)
        p_title_r.font.bold = True
        p_title_r.font.color.rgb = theme["accent"]
        p_title_r.space_after = Pt(12)
        
    if content.column2_points:
        for idx, pt in enumerate(content.column2_points):
            p = tf_r.add_paragraph() if (idx > 0 or content.column2_title) else tf_r.paragraphs[0]
            p.text = f"•  {pt}"
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = theme["secondary"]
            p.space_after = Pt(10)

def add_header(slide, title: str, theme: dict):
    """輔助函式：在投影片頂端加入工整的標題與下底線線條。"""
    tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["primary"]
    
    # 標題下方小裝飾線 (寬度 1.5 英吋，厚度 0.04 英吋)
    shape = slide.shapes.add_shape(
        1,  # msoShapeRectangle
        Inches(1.0), Inches(1.7), Inches(1.5), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["accent"]
    shape.line.fill.background()
